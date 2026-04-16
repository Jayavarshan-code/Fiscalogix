"""
Fiscalogix Pitch Deck Generator
Run: python generate_pitch_deck.py
Output: fiscalogix_pitch_deck.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── Colour Palette ──────────────────────────────────────────────────────────
NAVY       = RGBColor(0x0A, 0x16, 0x28)
BLUE_DARK  = RGBColor(0x1E, 0x3A, 0x8A)
BLUE_MID   = RGBColor(0x1D, 0x4E, 0xD8)
BLUE_LIGHT = RGBColor(0x93, 0xC5, 0xFD)
GOLD       = RGBColor(0xF5, 0x9E, 0x0B)
RED_ACC    = RGBColor(0xEF, 0x44, 0x44)
GREEN      = RGBColor(0x10, 0xB9, 0x81)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
GRAY_MID   = RGBColor(0x64, 0x74, 0x8B)
GRAY_DARK  = RGBColor(0x1E, 0x29, 0x3B)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def prs_new():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    blank = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(blank)


# ── Low-level shape helpers ─────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_color, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox


def add_text_multiline(slide, lines, left, top, width, height,
                        font_size=14, color=WHITE, bold=False, spacing=1.2):
    """lines = list of strings"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return txBox


def gradient_background(slide, color1, color2):
    """Add a full-slide gradient rectangle (approximated with two overlapping rects)."""
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, color1)
    # Semi-transparent overlay hint using a lighter rect on the right half
    overlay = add_rect(slide, Inches(6), 0, Inches(7.33), SLIDE_H, color2)
    overlay.fill.fore_color.rgb = color2


def accent_bar(slide, color=GOLD, top=None, height=Pt(4)):
    if top is None:
        top = SLIDE_H - height
    add_rect(slide, 0, top, SLIDE_W, int(height), color)


def section_badge(slide, text, left, top, bg=BLUE_MID, fg=WHITE):
    w, h = Inches(2.2), Inches(0.32)
    add_rect(slide, left, top, w, h, bg)
    add_text(slide, text.upper(), left + Inches(0.12), top + Inches(0.02),
             w - Inches(0.24), h, font_size=9, bold=True, color=fg,
             align=PP_ALIGN.LEFT)


def bullet_point(slide, text, left, top, width, dot_color=GOLD, font_size=13, color=WHITE):
    # dot
    add_rect(slide, left, top + Inches(0.07), Inches(0.07), Inches(0.07), dot_color)
    add_text(slide, text, left + Inches(0.2), top, width - Inches(0.2),
             Inches(0.3), font_size=font_size, color=color)


def card(slide, left, top, width, height, fill=GRAY_DARK, title=None, title_color=GOLD,
         title_size=12, body_lines=None, body_color=WHITE, body_size=11):
    add_rect(slide, left, top, width, height, fill)
    cy = top + Inches(0.18)
    if title:
        add_text(slide, title, left + Inches(0.2), cy, width - Inches(0.4),
                 Inches(0.35), font_size=title_size, bold=True, color=title_color)
        cy += Inches(0.35)
    if body_lines:
        for line in body_lines:
            add_text(slide, "  · " + line, left + Inches(0.1), cy,
                     width - Inches(0.2), Inches(0.28), font_size=body_size, color=body_color)
            cy += Inches(0.28)


# ════════════════════════════════════════════════════════════════════════════
# SLIDES
# ════════════════════════════════════════════════════════════════════════════

def slide_01_cover(prs):
    s = blank_slide(prs)
    # Background
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_rect(s, Inches(7.5), 0, Inches(5.83), SLIDE_H, BLUE_DARK)  # Right panel

    # Accent dots / decoration
    for i in range(6):
        add_rect(s, Inches(8 + i * 0.8), Inches(0.3), Inches(0.4), Inches(0.4),
                 RGBColor(0x1D + i*5, 0x4E, 0xD8))

    # Logo block
    add_rect(s, Inches(0.7), Inches(1.8), Inches(0.7), Inches(0.7), GOLD)
    add_text(s, "F", Inches(0.7), Inches(1.82), Inches(0.7), Inches(0.7),
             font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(s, "FISCALOGIX", Inches(1.55), Inches(1.9),
             Inches(5), Inches(0.6), font_size=36, bold=True, color=WHITE)
    add_text(s, "Intelligent Financial Operating System",
             Inches(1.55), Inches(2.55), Inches(6), Inches(0.4),
             font_size=14, color=BLUE_LIGHT, italic=True)

    # Divider
    add_rect(s, Inches(0.7), Inches(3.1), Inches(5), Inches(0.04), GOLD)

    # Sub-headline
    add_text_multiline(s,
        ["India's first AI-native platform unifying financial management,",
         "supply chain intelligence, GST compliance and demand forecasting."],
        Inches(0.7), Inches(3.3), Inches(6.2), Inches(1.0),
        font_size=13, color=RGBColor(0xCB, 0xD5, 0xE1))

    # Badges
    badges = ["AI-Powered", "GST Ready", "Real-Time", "Enterprise"]
    for i, b in enumerate(badges):
        bx = Inches(0.7 + i * 1.5)
        add_rect(s, bx, Inches(4.5), Inches(1.35), Inches(0.35), BLUE_MID)
        add_text(s, b, bx + Inches(0.1), Inches(4.52), Inches(1.15), Inches(0.3),
                 font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Right panel content
    add_text(s, "PITCH DECK", Inches(8.2), Inches(1.5), Inches(4), Inches(0.4),
             font_size=10, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(s, "2025", Inches(8.2), Inches(2.0), Inches(4), Inches(0.6),
             font_size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    stats = [("6+", "AI/ML Models"), ("90d", "Forecast Range"),
             ("1", "Unified Platform"), ("∞", "Audit Trails")]
    for i, (num, lbl) in enumerate(stats):
        bx = Inches(8.2 + (i % 2) * 2.2)
        by = Inches(3.0 + (i // 2) * 1.5)
        add_rect(s, bx, by, Inches(1.9), Inches(1.2), NAVY)
        add_text(s, num, bx, by + Inches(0.1), Inches(1.9), Inches(0.6),
                 font_size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_text(s, lbl, bx, by + Inches(0.7), Inches(1.9), Inches(0.4),
                 font_size=9, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)

    accent_bar(s, GOLD)
    accent_bar(s, RED_ACC, top=SLIDE_H - Pt(8), height=Pt(4))


def slide_02_problem(prs):
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, RGBColor(0xFD, 0xF2, 0xF2))
    add_rect(s, 0, 0, SLIDE_W, Inches(1.1), NAVY)
    accent_bar(s, RED_ACC, top=Inches(1.1) - Pt(4), height=Pt(4))

    add_text(s, "THE PROBLEM", Inches(0.5), Inches(0.25), Inches(8), Inches(0.5),
             font_size=26, bold=True, color=WHITE)
    add_text(s, "How Indian enterprises manage finances today — and why it's broken",
             Inches(0.5), Inches(0.75), Inches(10), Inches(0.3),
             font_size=12, color=BLUE_LIGHT, italic=True)

    problems = [
        ("Fragmented Tools", "Tally, Excel, GST portal, logistics sheets — none talk to each other. Decisions happen on stale data."),
        ("GST Nightmare", "Reconciliation takes finance teams days each month. Errors mean penalties. ITC mismatches go unnoticed."),
        ("Blind to Disruptions", "Supply chain disruptions — port strikes, geopolitical risk, weather — are noticed after the damage is done."),
        ("No Predictive Power", "Cash flow projections are monthly spreadsheets. Demand signals arrive too late to act on inventory."),
        ("Audit Chaos", "Audit trails are fragmented across 5 systems. Producing compliance evidence is a manual, multi-day exercise."),
    ]

    for i, (title, desc) in enumerate(problems):
        row = i // 2
        col = i % 2
        if i == 4:
            bx, by = Inches(3.5), Inches(1.4 + row * 1.65)
            bw = Inches(6.0)
        else:
            bx = Inches(0.4 + col * 6.4)
            by = Inches(1.4 + row * 1.65)
            bw = Inches(6.0)

        add_rect(s, bx, by, bw, Inches(1.45), WHITE)
        add_rect(s, bx, by, Inches(0.06), Inches(1.45), RED_ACC)
        add_text(s, "✗  " + title, bx + Inches(0.2), by + Inches(0.12), bw - Inches(0.3),
                 Inches(0.35), font_size=13, bold=True, color=RGBColor(0xDC, 0x26, 0x26))
        add_text(s, desc, bx + Inches(0.2), by + Inches(0.5), bw - Inches(0.3),
                 Inches(0.9), font_size=11, color=GRAY_DARK)

    accent_bar(s, RED_ACC)


def slide_03_solution(prs):
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_rect(s, 0, 0, SLIDE_W, Inches(1.1), BLUE_DARK)
    accent_bar(s, GOLD, top=Inches(1.1) - Pt(4), height=Pt(4))

    add_text(s, "THE SOLUTION", Inches(0.5), Inches(0.25), Inches(8), Inches(0.5),
             font_size=26, bold=True, color=WHITE)
    add_text(s, "One intelligent platform — built ground-up for Indian enterprises",
             Inches(0.5), Inches(0.75), Inches(10), Inches(0.3),
             font_size=12, color=BLUE_LIGHT, italic=True)

    solutions = [
        ("🧠 AI Digital Twin", "Live simulation of your entire operation. Ask 'what if' questions and get answers in seconds."),
        ("📊 Demand Forecasting", "90-day ML forecasts with SHAP explainability. Know what to stock before demand hits."),
        ("🛣️ Route Optimizer", "Geopolitical-aware routing across Ocean, Rail & Truck. Avoids strikes, sanctions & delays."),
        ("🇮🇳 GST Compliance", "Auto-generate GSTR-1, GSTR-3B. Reconcile ITC. Zero manual GST work."),
        ("⚡ Real-Time Alerts", "20+ financial risk signals monitored 24/7. SMS + dashboard alerts before crises hit."),
        ("💬 AI Copilot", "Claude-powered RAG assistant over your contracts, invoices & financial history."),
    ]

    for i, (title, desc) in enumerate(solutions):
        row = i // 3
        col = i % 3
        bx = Inches(0.4 + col * 4.3)
        by = Inches(1.5 + row * 2.5)
        bw, bh = Inches(4.0), Inches(2.2)

        add_rect(s, bx, by, bw, bh, BLUE_DARK)
        add_rect(s, bx, by, bw, Inches(0.05), GOLD)
        add_text(s, title, bx + Inches(0.2), by + Inches(0.15), bw - Inches(0.4),
                 Inches(0.45), font_size=13, bold=True, color=GOLD)
        add_text(s, desc, bx + Inches(0.2), by + Inches(0.65), bw - Inches(0.4),
                 Inches(1.4), font_size=11, color=RGBColor(0xCB, 0xD5, 0xE1))

    accent_bar(s, GOLD)


def slide_04_product(prs):
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, RGBColor(0xF8, 0xFA, 0xFF))
    add_rect(s, 0, 0, SLIDE_W, Inches(1.1), NAVY)
    add_rect(s, 0, 0, SLIDE_W, Inches(1.1), NAVY)
    accent_bar(s, BLUE_MID, top=Inches(1.1) - Pt(4), height=Pt(4))

    add_text(s, "HOW IT WORKS", Inches(0.5), Inches(0.25), Inches(8), Inches(0.5),
             font_size=26, bold=True, color=WHITE)
    add_text(s, "The Fiscalogix architecture — from raw data to intelligent action",
             Inches(0.5), Inches(0.75), Inches(10), Inches(0.3),
             font_size=12, color=BLUE_LIGHT, italic=True)

    # Flow diagram (simplified as boxes with arrows)
    layers = [
        ("DATA LAYER", ["ERP / Tally", "GST Portal", "Logistics APIs", "Bank Feeds", "Market Data"],
         RGBColor(0x0F, 0x17, 0x2A), RGBColor(0x1E, 0x2D, 0x4F), RGBColor(0x1A, 0x24, 0x3D)),
        ("AI ENGINE", ["Adaptive Orchestrator", "Risk Engine", "Demand Models", "Route Optimizer", "RAG Embeddings"],
         BLUE_DARK, BLUE_MID, RGBColor(0x1E, 0x50, 0xAA)),
        ("OUTPUT", ["Digital Twin Dashboard", "GST Reports", "Alerts & SMS", "Audit Logs", "AI Copilot"],
         RGBColor(0x05, 0x4F, 0x31), RGBColor(0x06, 0x78, 0x4C), RGBColor(0x06, 0x64, 0x40)),
    ]

    for i, (label, items, color, header_color, item_color) in enumerate(layers):
        bx = Inches(0.5 + i * 4.2)
        by = Inches(1.5)
        bw, bh = Inches(3.8), Inches(5.5)
        add_rect(s, bx, by, bw, bh, color)
        add_rect(s, bx, by, bw, Inches(0.45), header_color)
        add_text(s, label, bx + Inches(0.1), by + Inches(0.05), bw - Inches(0.2), Inches(0.35),
                 font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        for j, item in enumerate(items):
            iy = by + Inches(0.65 + j * 0.85)
            add_rect(s, bx + Inches(0.2), iy, bw - Inches(0.4), Inches(0.6), item_color)
            add_text(s, item, bx + Inches(0.3), iy + Inches(0.1),
                     bw - Inches(0.5), Inches(0.4), font_size=11, color=WHITE)

        # Arrow between boxes
        if i < 2:
            ax = bx + bw + Inches(0.05)
            ay = by + Inches(2.5)
            add_text(s, "→", ax, ay, Inches(0.3), Inches(0.4),
                     font_size=20, bold=True, color=BLUE_MID, align=PP_ALIGN.CENTER)

    accent_bar(s, BLUE_MID)


def slide_05_market(prs):
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_rect(s, 0, 0, SLIDE_W, Inches(1.1), BLUE_DARK)
    accent_bar(s, GOLD, top=Inches(1.1) - Pt(4), height=Pt(4))

    add_text(s, "MARKET OPPORTUNITY", Inches(0.5), Inches(0.25), Inches(8), Inches(0.5),
             font_size=26, bold=True, color=WHITE)
    add_text(s, "India's enterprise software market is massively underserved",
             Inches(0.5), Inches(0.75), Inches(10), Inches(0.3),
             font_size=12, color=BLUE_LIGHT, italic=True)

    circles = [
        ("TAM", "$12B", "India Enterprise\nFinancial Software\nMarket (2025)", Inches(1.5), BLUE_DARK),
        ("SAM", "$3.2B", "SME + Mid-Market\nGST + Supply Chain\nSoftware", Inches(4.3), BLUE_MID),
        ("SOM", "$320M", "Immediately\nAddressable with\nFiscalogix v1", Inches(7.1), GOLD),
    ]

    for label, amount, desc, cx, color in circles:
        is_gold = color == GOLD
        fg = NAVY if is_gold else WHITE
        add_rect(s, cx, Inches(1.7), Inches(2.5), Inches(2.5), color)
        add_text(s, label, cx, Inches(1.8), Inches(2.5), Inches(0.4),
                 font_size=10, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF) if not is_gold else NAVY,
                 align=PP_ALIGN.CENTER)
        add_text(s, amount, cx, Inches(2.25), Inches(2.5), Inches(0.7),
                 font_size=30, bold=True, color=fg, align=PP_ALIGN.CENTER)
        add_text(s, desc, cx, Inches(2.9), Inches(2.5), Inches(0.9),
                 font_size=9, color=RGBColor(0xCC, 0xCC, 0xCC) if not is_gold else NAVY,
                 align=PP_ALIGN.CENTER)

    # Why now bullets
    add_text(s, "Why Now", Inches(0.5), Inches(4.4), Inches(4), Inches(0.4),
             font_size=14, bold=True, color=GOLD)
    why_now = [
        "GST e-invoicing mandate expands to all businesses in 2025",
        "India's logistics sector grows 15% YoY — needs intelligent routing",
        "MSME digitisation push: ₹22,000 Cr government backing",
        "No incumbent offers AI + GST + Supply Chain in one platform",
    ]
    for i, pt in enumerate(why_now):
        add_rect(s, Inches(0.5), Inches(4.9 + i * 0.42), Inches(0.08), Inches(0.25), GOLD)
        add_text(s, pt, Inches(0.75), Inches(4.87 + i * 0.42), Inches(11), Inches(0.35),
                 font_size=11, color=RGBColor(0xCB, 0xD5, 0xE1))

    accent_bar(s, GOLD)


def slide_06_business_model(prs):
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, RGBColor(0xF8, 0xFA, 0xFF))
    add_rect(s, 0, 0, SLIDE_W, Inches(1.1), NAVY)
    accent_bar(s, GREEN, top=Inches(1.1) - Pt(4), height=Pt(4))

    add_text(s, "BUSINESS MODEL", Inches(0.5), Inches(0.25), Inches(8), Inches(0.5),
             font_size=26, bold=True, color=WHITE)
    add_text(s, "SaaS + API — recurring revenue with usage-based upsell",
             Inches(0.5), Inches(0.75), Inches(10), Inches(0.3),
             font_size=12, color=BLUE_LIGHT, italic=True)

    tiers = [
        ("STARTER", "₹9,999/mo", GRAY_DARK, WHITE,
         ["GST Compliance Module", "Basic Demand Forecasting", "Dashboard & Alerts", "Up to 5 users", "Email support"]),
        ("GROWTH", "₹29,999/mo", BLUE_DARK, WHITE,
         ["Full Digital Twin", "Route Optimizer", "AI Copilot (RAG)", "Up to 25 users", "API Access + Webhooks"]),
        ("ENTERPRISE", "Custom", NAVY, GOLD,
         ["On-premise / Private Cloud", "Custom ML fine-tuning", "Dedicated SLA", "Unlimited users", "Full integrations"]),
    ]

    for i, (tier, price, bg, title_color, features) in enumerate(tiers):
        bx = Inches(0.5 + i * 4.2)
        bw = Inches(3.9)
        add_rect(s, bx, Inches(1.4), bw, Inches(5.6), bg)
        if tier == "GROWTH":
            add_rect(s, bx, Inches(1.4), bw, Inches(0.08), GREEN)
        add_text(s, tier, bx + Inches(0.2), Inches(1.5), bw - Inches(0.4), Inches(0.4),
                 font_size=11, bold=True, color=GOLD if tier != "GROWTH" else GREEN,
                 align=PP_ALIGN.CENTER)
        add_text(s, price, bx + Inches(0.1), Inches(2.0), bw - Inches(0.2), Inches(0.6),
                 font_size=26, bold=True, color=title_color, align=PP_ALIGN.CENTER)
        add_rect(s, bx + Inches(0.5), Inches(2.65), bw - Inches(1.0), Inches(0.03),
                 RGBColor(0x33, 0x44, 0x55))
        for j, feat in enumerate(features):
            fy = Inches(2.8 + j * 0.6)
            add_rect(s, bx + Inches(0.3), fy + Inches(0.08), Inches(0.1), Inches(0.1),
                     GREEN if tier == "GROWTH" else GOLD)
            add_text(s, feat, bx + Inches(0.55), fy, bw - Inches(0.7), Inches(0.5),
                     font_size=11, color=WHITE if tier != "ENTERPRISE" else RGBColor(0xCB, 0xD5, 0xE1))

    # Revenue model note
    add_text(s, "Additional revenue: API usage metering · Professional services · Data analytics add-ons",
             Inches(0.5), Inches(7.1), Inches(12), Inches(0.3),
             font_size=10, color=GRAY_MID, align=PP_ALIGN.CENTER)

    accent_bar(s, GREEN)


def slide_07_traction(prs):
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_rect(s, 0, 0, SLIDE_W, Inches(1.1), BLUE_DARK)
    accent_bar(s, GOLD, top=Inches(1.1) - Pt(4), height=Pt(4))

    add_text(s, "TRACTION & ROADMAP", Inches(0.5), Inches(0.25), Inches(8), Inches(0.5),
             font_size=26, bold=True, color=WHITE)
    add_text(s, "Where we are and where we're going",
             Inches(0.5), Inches(0.75), Inches(10), Inches(0.3),
             font_size=12, color=BLUE_LIGHT, italic=True)

    # Left: what's built
    add_text(s, "Built & Deployed", Inches(0.5), Inches(1.3), Inches(5), Inches(0.4),
             font_size=14, bold=True, color=GOLD)
    built = [
        "Production backend on Railway (FastAPI + PostgreSQL)",
        "6 ML models: demand, risk, routing, RAG, anomaly detection",
        "GST compliance engine with GSTR-1 / GSTR-3B generation",
        "Digital twin with real-time simulation engine",
        "AI Copilot powered by Claude (Anthropic)",
        "SMS alerting via Twilio, async job processing via Celery",
    ]
    for i, item in enumerate(built):
        add_rect(s, Inches(0.5), Inches(1.85 + i * 0.58), Inches(0.1), Inches(0.25), GREEN)
        add_text(s, item, Inches(0.75), Inches(1.82 + i * 0.58), Inches(5.5), Inches(0.5),
                 font_size=11, color=RGBColor(0xCB, 0xD5, 0xE1))

    # Right: roadmap
    add_rect(s, Inches(7.0), Inches(1.3), Inches(5.8), Inches(5.8), BLUE_DARK)
    add_text(s, "Roadmap", Inches(7.2), Inches(1.45), Inches(5), Inches(0.4),
             font_size=14, bold=True, color=GOLD)

    roadmap = [
        ("Q2 2025", "First 5 paying pilot customers (SME segment)"),
        ("Q3 2025", "WhatsApp & ERP integration (SAP, Zoho)"),
        ("Q4 2025", "Custom ML fine-tuning on customer data"),
        ("Q1 2026", "Series A raise + 50 enterprise accounts"),
        ("Q2 2026", "Expand to SEA markets (Singapore, Indonesia)"),
    ]
    for i, (quarter, milestone) in enumerate(roadmap):
        by = Inches(2.0 + i * 0.95)
        add_rect(s, Inches(7.2), by, Inches(1.1), Inches(0.35), GOLD)
        add_text(s, quarter, Inches(7.2), by + Inches(0.02), Inches(1.1), Inches(0.3),
                 font_size=9, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(s, milestone, Inches(8.5), by, Inches(4.0), Inches(0.5),
                 font_size=11, color=WHITE)

    accent_bar(s, GOLD)


def slide_08_team(prs):
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, RGBColor(0xF8, 0xFA, 0xFF))
    add_rect(s, 0, 0, SLIDE_W, Inches(1.1), NAVY)
    accent_bar(s, BLUE_MID, top=Inches(1.1) - Pt(4), height=Pt(4))

    add_text(s, "THE TEAM", Inches(0.5), Inches(0.25), Inches(8), Inches(0.5),
             font_size=26, bold=True, color=WHITE)
    add_text(s, "Built by engineers who understand both finance and AI",
             Inches(0.5), Inches(0.75), Inches(10), Inches(0.3),
             font_size=12, color=BLUE_LIGHT, italic=True)

    # Founder card (centered, prominent)
    add_rect(s, Inches(1.5), Inches(1.5), Inches(10.3), Inches(2.0), NAVY)
    add_rect(s, Inches(1.5), Inches(1.5), Inches(0.07), Inches(2.0), GOLD)
    add_text(s, "Varshan", Inches(1.8), Inches(1.65), Inches(4), Inches(0.5),
             font_size=20, bold=True, color=GOLD)
    add_text(s, "Founder & CEO — Fiscalogix", Inches(1.8), Inches(2.15), Inches(4), Inches(0.35),
             font_size=12, color=BLUE_LIGHT)
    add_text(s, "Full-stack AI engineer with deep expertise in financial systems, ML pipelines,\nand enterprise SaaS architecture. Built Fiscalogix end-to-end.",
             Inches(1.8), Inches(2.55), Inches(9.5), Inches(0.7),
             font_size=11, color=RGBColor(0xCB, 0xD5, 0xE1))

    # What we're looking for
    add_text(s, "We Are Looking For", Inches(0.5), Inches(3.9), Inches(5), Inches(0.4),
             font_size=14, bold=True, color=NAVY)

    looking = [
        ("Strategic Advisors", "CFOs, GST experts, logistics heads who can open doors and validate our roadmap"),
        ("Enterprise Pilots", "5–10 companies willing to be design partners in exchange for deeply customized onboarding"),
        ("Seed Investors", "₹2–5 Cr to extend runway to Series A, accelerate BD, and hire 2 senior engineers"),
    ]
    for i, (title, desc) in enumerate(looking):
        bx = Inches(0.5 + i * 4.2)
        add_rect(s, bx, Inches(4.5), Inches(3.9), Inches(2.4), BLUE_DARK)
        add_rect(s, bx, Inches(4.5), Inches(3.9), Inches(0.06), GOLD)
        add_text(s, title, bx + Inches(0.2), Inches(4.65), Inches(3.5), Inches(0.4),
                 font_size=12, bold=True, color=GOLD)
        add_text(s, desc, bx + Inches(0.2), Inches(5.15), Inches(3.5), Inches(1.5),
                 font_size=11, color=RGBColor(0xCB, 0xD5, 0xE1))

    accent_bar(s, BLUE_MID)


def slide_09_cta(prs):
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)

    # Bold centre gradient box
    add_rect(s, Inches(1.5), Inches(1.2), Inches(10.33), Inches(5.1), BLUE_DARK)
    add_rect(s, Inches(1.5), Inches(1.2), Inches(10.33), Inches(0.08), GOLD)
    add_rect(s, Inches(1.5), Inches(6.22), Inches(10.33), Inches(0.08), GOLD)

    add_text(s, "FISCALOGIX", Inches(1.5), Inches(1.6), Inches(10.33), Inches(0.8),
             font_size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "India's Intelligent Financial Operating System",
             Inches(1.5), Inches(2.45), Inches(10.33), Inches(0.5),
             font_size=16, color=BLUE_LIGHT, align=PP_ALIGN.CENTER, italic=True)

    add_rect(s, Inches(5.0), Inches(3.0), Inches(3.33), Inches(0.05), GOLD)

    add_text(s, "Let's Build the Future of Indian Enterprise Finance — Together",
             Inches(1.5), Inches(3.2), Inches(10.33), Inches(0.7),
             font_size=14, color=RGBColor(0xCB, 0xD5, 0xE1), align=PP_ALIGN.CENTER)

    ctas = ["Schedule a Demo", "Become a Pilot Partner", "Invest in Fiscalogix"]
    for i, cta in enumerate(ctas):
        bx = Inches(1.8 + i * 3.3)
        add_rect(s, bx, Inches(4.1), Inches(2.9), Inches(0.55),
                 GOLD if i == 0 else (BLUE_MID if i == 1 else GREEN))
        add_text(s, cta, bx, Inches(4.12), Inches(2.9), Inches(0.5),
                 font_size=12, bold=True,
                 color=NAVY if i == 0 else WHITE, align=PP_ALIGN.CENTER)

    add_text(s, "Confidential — For Discussion Purposes Only · Fiscalogix 2025",
             Inches(1.5), Inches(5.9), Inches(10.33), Inches(0.3),
             font_size=9, color=GRAY_MID, align=PP_ALIGN.CENTER)

    accent_bar(s, GOLD)
    accent_bar(s, RED_ACC, top=SLIDE_H - Pt(8), height=Pt(4))


# ════════════════════════════════════════════════════════════════════════════
# MAIN
# ════════════════════════════════════════════════════════════════════════════

def build():
    prs = prs_new()
    slide_01_cover(prs)
    slide_02_problem(prs)
    slide_03_solution(prs)
    slide_04_product(prs)
    slide_05_market(prs)
    slide_06_business_model(prs)
    slide_07_traction(prs)
    slide_08_team(prs)
    slide_09_cta(prs)

    out = "fiscalogix_pitch_deck.pptx"
    prs.save(out)
    print(f"[OK] Saved: {out}  ({len(prs.slides)} slides)")
    return out


if __name__ == "__main__":
    build()
