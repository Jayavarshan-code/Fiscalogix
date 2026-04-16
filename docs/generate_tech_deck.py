"""
Fiscalogix — Technical Capabilities Pitch Deck
Run: python generate_tech_deck.py
Output: fiscalogix_tech_deck.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colour Palette ──────────────────────────────────────────────────────────
NAVY        = RGBColor(0x07, 0x10, 0x24)
BLUE_DARK   = RGBColor(0x0F, 0x2A, 0x5E)
BLUE_MID    = RGBColor(0x1D, 0x4E, 0xD8)
BLUE_LIGHT  = RGBColor(0x93, 0xC5, 0xFD)
CYAN        = RGBColor(0x06, 0xB6, 0xD4)
GOLD        = RGBColor(0xF5, 0x9E, 0x0B)
AMBER       = RGBColor(0xFB, 0xBF, 0x24)
RED_ACC     = RGBColor(0xEF, 0x44, 0x44)
GREEN       = RGBColor(0x10, 0xB9, 0x81)
PURPLE      = RGBColor(0x7C, 0x3A, 0xED)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_LIGHT  = RGBColor(0xF1, 0xF5, 0xF9)
GRAY_MID    = RGBColor(0x64, 0x74, 0x8B)
GRAY_DARK   = RGBColor(0x1E, 0x29, 0x3B)
PANEL_BG    = RGBColor(0x0C, 0x1A, 0x36)
CARD_BG     = RGBColor(0x11, 0x22, 0x44)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helpers ─────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, l, t, w, h, color):
    sp = slide.shapes.add_shape(1, l, t, w, h)
    sp.line.fill.background()
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    return sp


def txt(slide, text, l, t, w, h, size=12, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = color
    run.font.name    = "Calibri"
    return box


def header_band(slide, title, subtitle=None, accent=GOLD):
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    rect(slide, 0, 0, SLIDE_W, Inches(1.05), PANEL_BG)
    rect(slide, 0, Inches(1.05) - Pt(4), SLIDE_W, Pt(4), accent)
    txt(slide, title, Inches(0.5), Inches(0.2), Inches(10), Inches(0.55),
        size=24, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, Inches(0.5), Inches(0.7), Inches(11), Inches(0.3),
            size=11, color=BLUE_LIGHT, italic=True)


def bottom_bar(slide, color=GOLD):
    rect(slide, 0, SLIDE_H - Pt(5), SLIDE_W, Pt(5), color)


def pill(slide, label, l, t, bg=BLUE_MID, fg=WHITE, size=9):
    w, h = Inches(1.5), Inches(0.28)
    rect(slide, l, t, w, h, bg)
    txt(slide, label, l + Inches(0.06), t + Inches(0.02),
        w - Inches(0.12), h - Inches(0.04), size=size, bold=True,
        color=fg, align=PP_ALIGN.CENTER)


def card(slide, l, t, w, h, bg=CARD_BG, top_accent=None):
    rect(slide, l, t, w, h, bg)
    if top_accent:
        rect(slide, l, t, w, Pt(4), top_accent)


def dot_bullet(slide, text, l, t, w, dot=GOLD, size=11, color=WHITE):
    rect(slide, l, t + Inches(0.08), Inches(0.07), Inches(0.07), dot)
    txt(slide, text, l + Inches(0.18), t, w - Inches(0.18), Inches(0.35),
        size=size, color=color)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ════════════════════════════════════════════════════════════════════════════
def s01_cover(prs):
    s = blank(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    # Left accent strip
    rect(s, 0, 0, Inches(0.12), SLIDE_H, GOLD)
    # Top bar
    rect(s, 0, 0, SLIDE_W, Inches(0.08), GOLD)
    # Right dark panel
    rect(s, Inches(8.2), 0, Inches(5.13), SLIDE_H, PANEL_BG)

    # Logo block
    rect(s, Inches(0.5), Inches(1.6), Inches(0.7), Inches(0.7), GOLD)
    txt(s, "F", Inches(0.5), Inches(1.62), Inches(0.7), Inches(0.7),
        size=30, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    txt(s, "FISCALOGIX", Inches(1.35), Inches(1.72), Inches(6.5), Inches(0.65),
        size=38, bold=True, color=WHITE)
    txt(s, "TECHNICAL CAPABILITIES BRIEF",
        Inches(1.35), Inches(2.42), Inches(6.5), Inches(0.38),
        size=13, bold=True, color=GOLD)

    rect(s, Inches(0.5), Inches(3.0), Inches(5.5), Inches(0.04), BLUE_MID)

    txt(s, "A full-stack AI financial operating system built for enterprises\n"
           "that run complex supply chains, manage Indian tax compliance,\n"
           "and need predictive intelligence — not just reporting.",
        Inches(0.5), Inches(3.15), Inches(7.2), Inches(1.1),
        size=12, color=RGBColor(0xCB, 0xD5, 0xE1))

    pills = [("8-Stage AI Pipeline", BLUE_MID), ("GNN Risk Engine", PURPLE),
             ("GST Compliance", GREEN), ("Monte Carlo", CYAN)]
    for i, (lbl, clr) in enumerate(pills):
        pill(s, lbl, Inches(0.5 + i * 1.75), Inches(4.4), bg=clr)

    # Right panel — stat grid
    txt(s, "WHAT'S INSIDE", Inches(8.5), Inches(0.9), Inches(4.5), Inches(0.35),
        size=9, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    stats = [
        ("8", "Pipeline Stages"),
        ("6+", "ML Models"),
        ("20+", "Risk Signals"),
        ("5", "AI Agents"),
        ("RAG", "Semantic Retrieval"),
        ("GSTR", "Auto Compliance"),
    ]
    for i, (val, lbl) in enumerate(stats):
        col, row = i % 2, i // 2
        bx = Inches(8.4 + col * 2.3)
        by = Inches(1.4 + row * 1.6)
        card(s, bx, by, Inches(2.1), Inches(1.35), CARD_BG, GOLD if col == 0 else BLUE_MID)
        txt(s, val, bx, by + Inches(0.15), Inches(2.1), Inches(0.6),
            size=30, bold=True, color=GOLD if col == 0 else CYAN, align=PP_ALIGN.CENTER)
        txt(s, lbl, bx, by + Inches(0.75), Inches(2.1), Inches(0.45),
            size=9, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)

    bottom_bar(s, GOLD)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — WHY IT'S NEEDED
# ════════════════════════════════════════════════════════════════════════════
def s02_why(prs):
    s = blank(prs)
    header_band(s, "WHY THIS PRODUCT IS NEEDED",
                "The exact failures that led to building Fiscalogix from scratch", RED_ACC)

    problems = [
        ("No unified financial intelligence",
         "Companies run Tally for accounts, Excel for forecasting, a third-party GST tool, and a logistics spreadsheet. "
         "None share data. Decisions happen on stale, disconnected snapshots.",
         RED_ACC),
        ("GST ITC lag is an invisible working capital drain",
         "Importing goods to India triggers IGST at customs. Recovery as Input Tax Credit takes 45-90 days. "
         "At 11% WACC, a Rs1Cr import bleeds Rs24,000 per shipment in hidden carrying cost — invisible in any P&L.",
         AMBER),
        ("Risk is detected AFTER the damage",
         "Supply chain disruptions — port strikes, geopolitical closures, carrier defaults — surface in morning "
         "reports after shipments are already delayed. No system monitors leading indicators in real time.",
         RED_ACC),
        ("Monte Carlo simulations use garbage inputs",
         "Existing tools run scenarios with uncapped Pareto distributions — producing outputs like "
         "'$500M VaR on a $10K shipment'. CFOs dismiss the models entirely. Capital reserves are guesswork.",
         AMBER),
        ("ML models embedded in wrong industry assumptions",
         "Off-the-shelf demand models apply B2C retail seasonality (Q4 peak) to B2B manufacturers. "
         "An automotive forwarder gets told to over-stock in December when their peak is January factory restart.",
         RED_ACC),
        ("No audit trail across the decision chain",
         "When a compliance auditor asks 'why did the system route this shipment through Singapore?', "
         "no answer exists. Decisions are made, executed and forgotten. Accountability is impossible.",
         AMBER),
    ]

    for i, (title, desc, accent) in enumerate(problems):
        row, col = i // 2, i % 2
        bx = Inches(0.35 + col * 6.5)
        by = Inches(1.2 + row * 1.98)
        card(s, bx, by, Inches(6.2), Inches(1.78), CARD_BG, accent)
        txt(s, title, bx + Inches(0.2), by + Inches(0.18), Inches(5.8), Inches(0.38),
            size=11, bold=True, color=accent)
        txt(s, desc, bx + Inches(0.2), by + Inches(0.62), Inches(5.8), Inches(1.05),
            size=10, color=RGBColor(0xCB, 0xD5, 0xE1))

    bottom_bar(s, RED_ACC)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — SYSTEM ARCHITECTURE (8-Stage Pipeline)
# ════════════════════════════════════════════════════════════════════════════
def s03_architecture(prs):
    s = blank(prs)
    header_band(s, "8-STAGE AI PIPELINE ARCHITECTURE",
                "AdaptiveOrchestrator — replaces a single 80-line God-function with isolated, fault-tolerant stages", CYAN)

    stages = [
        ("1", "Data Ingestion", "Load + normalize from ERP,\nCSV, GST portal, bank feeds", BLUE_MID),
        ("2", "ML Inference", "Delay prediction + demand\nbatch inference (RandomForest)", PURPLE),
        ("3", "CLV Calibration", "Per-account Customer Lifetime\nValue enrichment", CYAN),
        ("4", "Decision Engine", "Per-row deterministic rules\n+ ReVM scoring", GOLD),
        ("5", "Situation Assessment", "Portfolio heuristics\n< 1ms, no LLM call", GREEN),
        ("6", "Dispatch Planning", "LLM agent selection\ntemp=0, deterministic", AMBER),
        ("7", "Agent Execution", "Runs Risk / Financial /\nRouting / Anomaly agents", RED_ACC),
        ("8", "Persistence", "Audit log + snapshot +\ndecision log written", BLUE_LIGHT),
    ]

    # Arrow rail
    rail_y = Inches(3.6)
    rect(s, Inches(0.3), rail_y + Inches(0.15), Inches(12.73), Inches(0.06), BLUE_DARK)

    for i, (num, title, desc, color) in enumerate(stages):
        bx = Inches(0.25 + i * 1.6)
        # connector dot on rail
        rect(s, bx + Inches(0.52), rail_y + Inches(0.08), Inches(0.2), Inches(0.2), color)
        # upper card
        card(s, bx, Inches(1.3), Inches(1.5), Inches(2.1), CARD_BG, color)
        txt(s, num, bx, Inches(1.35), Inches(1.5), Inches(0.4),
            size=18, bold=True, color=color, align=PP_ALIGN.CENTER)
        txt(s, title, bx + Inches(0.06), Inches(1.78), Inches(1.38), Inches(0.38),
            size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # connector line down
        rect(s, bx + Inches(0.72), Inches(3.4), Inches(0.06), Inches(0.22), color)
        # lower desc card
        card(s, bx, Inches(3.85), Inches(1.5), Inches(1.35), PANEL_BG, color)
        txt(s, desc, bx + Inches(0.08), Inches(3.95), Inches(1.34), Inches(1.15),
            size=8, color=RGBColor(0xCB, 0xD5, 0xE1))

    # Key guarantee box
    card(s, Inches(0.3), Inches(5.4), Inches(12.73), Inches(1.7), PANEL_BG, GREEN)
    txt(s, "FAULT-TOLERANCE GUARANTEE", Inches(0.5), Inches(5.5), Inches(6), Inches(0.35),
        size=10, bold=True, color=GREEN)
    guarantees = [
        "Stage N failure is caught, logged, and stored as StageOutput. Stage N+1 always executes.",
        "ctx.result('stage_name') returns {} safely — downstream stages never see a raw exception.",
        "pipeline_health key in every response exposes per-stage timings and failure list for observability.",
        "GNN failure cannot silently corrupt the ReVM calculation (isolated stage boundary).",
    ]
    for i, g in enumerate(guarantees):
        dot_bullet(s, g, Inches(0.5 + (i % 2) * 6.3), Inches(5.95 + (i // 2) * 0.42),
                   Inches(5.8), dot=GREEN, size=10)

    bottom_bar(s, CYAN)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — RISK ENGINE (XGBoost + GNN)
# ════════════════════════════════════════════════════════════════════════════
def s04_risk(prs):
    s = blank(prs)
    header_band(s, "RISK ENGINE — XGBoost + Graph Neural Network",
                "Two-model ensemble: tabular XGBoost + GNN contagion propagation across the supply graph", RED_ACC)

    # Left — XGBoost block
    card(s, Inches(0.3), Inches(1.2), Inches(5.8), Inches(5.8), CARD_BG, RED_ACC)
    txt(s, "XGBoost Risk Classifier", Inches(0.5), Inches(1.35), Inches(5.4), Inches(0.4),
        size=13, bold=True, color=RED_ACC)
    xgb_items = [
        "Trained on: carrier performance, route history, geopolitical scores, weather, SLA breach records",
        "SHAP explainability — every prediction has feature importances exposed to the frontend",
        "risk_pipeline.pkl: sklearn Pipeline (preprocessor + XGBoost classifier)",
        "Outputs: risk_probability, risk_class (LOW/MED/HIGH), top_shap_features",
        "Fallback: heuristic rule set if pkl missing (graceful degradation on first boot)",
    ]
    for i, item in enumerate(xgb_items):
        dot_bullet(s, item, Inches(0.5), Inches(1.9 + i * 0.6), Inches(5.4), size=10)

    txt(s, "Why SHAP matters to auditors:", Inches(0.5), Inches(4.95), Inches(5.4), Inches(0.3),
        size=10, bold=True, color=AMBER)
    txt(s, "Every risk decision can be traced to its input features.\n"
           "No black box — auditors see exactly why a shipment was flagged.",
        Inches(0.5), Inches(5.3), Inches(5.4), Inches(0.6), size=10, color=BLUE_LIGHT)

    # Right — GNN block
    card(s, Inches(6.4), Inches(1.2), Inches(6.63), Inches(5.8), CARD_BG, PURPLE)
    txt(s, "Graph Neural Network (Contagion Layer)", Inches(6.6), Inches(1.35),
        Inches(6.2), Inches(0.4), size=13, bold=True, color=PURPLE)
    gnn_items = [
        "Nodes = ports, warehouses, carriers. Edges = active shipping lanes",
        "Model: RiskGNN (PyTorch) — in_channels=2, hidden=16, out=2",
        "Predicts risk contagion: if Port A is disrupted, how does risk spread to Port B?",
        "GNN blend weight: up to 35% of final risk score when node is in path",
        "Decays with graph-hop distance — nodes 3+ hops away don't falsely alarm",
        "Temporal risk sensing: ETA-aware scoring (risk at arrival time, not today)",
    ]
    for i, item in enumerate(gnn_items):
        dot_bullet(s, item, Inches(6.6), Inches(1.9 + i * 0.6), Inches(6.1), dot=PURPLE, size=10)

    txt(s, "Blend formula:", Inches(6.6), Inches(5.0), Inches(6.1), Inches(0.3),
        size=10, bold=True, color=AMBER)
    card(s, Inches(6.6), Inches(5.35), Inches(6.1), Inches(0.55), NAVY)
    txt(s, "final_score = xgb_score * (1 - gnn_weight) + gnn_score * gnn_weight",
        Inches(6.7), Inches(5.42), Inches(5.9), Inches(0.4),
        size=10, color=CYAN)

    bottom_bar(s, RED_ACC)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — GST COMPLIANCE ENGINE
# ════════════════════════════════════════════════════════════════════════════
def s05_gst(prs):
    s = blank(prs)
    header_band(s, "INDIA GST COMPLIANCE ENGINE",
                "The only financial system that quantifies the working capital cost of GST on every shipment", GREEN)

    # Top row — 3 computation blocks
    blocks = [
        ("IMPORT (INBOUND)", [
            "IGST computed at customs clearance",
            "ITC recovery lag: 45-90 days modelled",
            "WACC-adjusted carrying cost calculated",
            "BCD (Basic Customs Duty) — sunk cost",
            "total_india_customs_cost returned per shipment",
        ], BLUE_MID, "Inbound to India"),
        ("EXPORT (OUTBOUND)", [
            "LUT path: zero upfront GST, zero working capital impact",
            "IGST-first path: refund lag cost modelled (30-90 days)",
            "DGFT Duty Drawback (AIR): % of FOB value recovered",
            "net_gst_impact = refund_lag_cost - drawback_benefit",
            "Flags exporters not on LUT for corrective action",
        ], GREEN, "Outbound from India"),
        ("GSTR AUTO-GENERATION", [
            "GSTR-1: outward supply return auto-populated",
            "GSTR-3B: summary return with ITC reconciliation",
            "Invoice validation against GST portal rules",
            "Mismatch detection before filing (not after penalty)",
            "Refund tracker: monitors GSTN processing status",
        ], GOLD, "Filing & Reconciliation"),
    ]

    for i, (title, items, color, label) in enumerate(blocks):
        bx = Inches(0.3 + i * 4.35)
        card(s, bx, Inches(1.2), Inches(4.1), Inches(4.3), CARD_BG, color)
        txt(s, title, bx + Inches(0.15), Inches(1.32), Inches(3.8), Inches(0.38),
            size=11, bold=True, color=color)
        txt(s, label, bx + Inches(0.15), Inches(1.72), Inches(3.8), Inches(0.28),
            size=8, color=GRAY_MID, italic=True)
        for j, item in enumerate(items):
            dot_bullet(s, item, bx + Inches(0.15), Inches(2.1 + j * 0.55),
                       Inches(3.75), dot=color, size=10)

    # Bottom — the killer insight
    card(s, Inches(0.3), Inches(5.65), Inches(12.73), Inches(1.55), PANEL_BG, AMBER)
    txt(s, "THE INSIGHT THAT NO OTHER TOOL COMPUTES", Inches(0.5), Inches(5.75),
        Inches(8), Inches(0.35), size=10, bold=True, color=AMBER)
    txt(s,
        "Rs 1 Cr pharma API import  |  IGST @ 12% = Rs 12L paid at port  |  ITC recovery: 67 days  |  WACC: 11%\n"
        "Hidden carrying cost = Rs 12L x 11% x (67/365) = Rs 24,200  per shipment  — completely invisible in any P&L.\n"
        "Fiscalogix surfaces this per-shipment, per-month. At scale: Rs 50L+ annual drain made visible and optimizable.",
        Inches(0.5), Inches(6.15), Inches(12.3), Inches(0.95), size=10, color=WHITE)

    bottom_bar(s, GREEN)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — GEOPOLITICAL ROUTE OPTIMIZER
# ════════════════════════════════════════════════════════════════════════════
def s06_route(prs):
    s = blank(prs)
    header_band(s, "GEOPOLITICAL ROUTE OPTIMIZER",
                "NetworkX-based multi-modal routing engine with real-time risk overlay and strike detection", CYAN)

    # Left: how it works
    card(s, Inches(0.3), Inches(1.2), Inches(6.1), Inches(5.8), CARD_BG, CYAN)
    txt(s, "How the Engine Works", Inches(0.5), Inches(1.35), Inches(5.7), Inches(0.38),
        size=13, bold=True, color=CYAN)

    how_items = [
        ("Graph Model", "Directed graph (NetworkX DiGraph). Nodes = ports/hubs/warehouses. Edges = shipping lanes."),
        ("Node Properties", "territory_type (Friendly/Neutral/Enemy), risk_score per node, ETA-aware scoring."),
        ("Edge Properties", "distance_km, duration_hours, base_cost (fuel+crew+fixed), transport_mode, customs_delay."),
        ("Transport Modes", "Ocean (25 km/h), Rail (60 km/h), Truck (80 km/h). Mode-switch penalty: $200 + 4 hrs."),
        ("Strike Detection", "set_strike(u, v) flags edges as active. Critical cargo = 2x penalty; standard = 0.5x."),
        ("Weight Function", "total_weight = base_cost x (1 + (node_risk + strike_penalty) x beta)"),
        ("Critical Cargo", "beta multiplied by 3.0 for critical shipments — routes aggressively away from risk."),
        ("Output", "Best path, feasibility score (0-100), risk level, cost breakdown by mode, carrier list."),
    ]
    for i, (label, desc) in enumerate(how_items):
        by = Inches(1.85 + i * 0.57)
        txt(s, label + ":", Inches(0.5), by, Inches(1.5), Inches(0.45),
            size=9, bold=True, color=GOLD)
        txt(s, desc, Inches(2.05), by, Inches(4.2), Inches(0.45), size=9,
            color=RGBColor(0xCB, 0xD5, 0xE1))

    # Right: why it matters
    card(s, Inches(6.7), Inches(1.2), Inches(6.33), Inches(2.5), CARD_BG, GOLD)
    txt(s, "Feasibility Score Formula", Inches(6.9), Inches(1.35), Inches(5.9), Inches(0.38),
        size=13, bold=True, color=GOLD)
    card(s, Inches(6.9), Inches(1.8), Inches(5.9), Inches(0.75), NAVY)
    txt(s, "feasibility = 100 - (mode_switches x 5) - (customs_delay x 2)\nclamped to [10, 98]",
        Inches(7.0), Inches(1.85), Inches(5.7), Inches(0.6), size=10, color=CYAN)
    score_items = [
        ("feasibility >= 85", "LOW risk — execute as planned", GREEN),
        ("feasibility 70-84", "MEDIUM risk — monitor closely", AMBER),
        ("feasibility < 70", "HIGH risk — consider rerouting", RED_ACC),
    ]
    for i, (cond, label, color) in enumerate(score_items):
        rect(s, Inches(6.9), Inches(2.7 + i * 0.27), Inches(1.5), Inches(0.22), color)
        txt(s, cond, Inches(6.92), Inches(2.71 + i * 0.27), Inches(1.46), Inches(0.2),
            size=8, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        txt(s, label, Inches(8.55), Inches(2.7 + i * 0.27), Inches(4.3), Inches(0.25),
            size=9, color=WHITE)

    card(s, Inches(6.7), Inches(3.9), Inches(6.33), Inches(3.1), CARD_BG, PURPLE)
    txt(s, "Multi-Echelon Inventory (MEIO)", Inches(6.9), Inches(4.05),
        Inches(5.9), Inches(0.38), size=13, bold=True, color=PURPLE)
    meio_items = [
        "Simulated Q-learning allocates stock across 3 hubs: NA (45%), EU (35%), APAC (remainder)",
        "WACC-adjusted holding cost = physical storage + capital carrying cost (opportunity cost)",
        "alpha = stockout_penalty / (true_holding_cost + stockout_penalty)",
        "Previous bug: WACC was accepted but never used — fixed, bias toward over-ordering eliminated",
        "projected_stockout_risk surfaced per SKU to decision engine",
    ]
    for i, item in enumerate(meio_items):
        dot_bullet(s, item, Inches(6.9), Inches(4.55 + i * 0.47),
                   Inches(5.9), dot=PURPLE, size=10)

    bottom_bar(s, CYAN)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — FINANCIAL INTELLIGENCE (ReVM + Monte Carlo)
# ════════════════════════════════════════════════════════════════════════════
def s07_financial(prs):
    s = blank(prs)
    header_band(s, "FINANCIAL INTELLIGENCE LAYER",
                "ReVM scoring, Monte Carlo VaR, cashflow prediction, liquidity scoring — all in one engine", GOLD)

    # ReVM formula
    card(s, Inches(0.3), Inches(1.2), Inches(12.73), Inches(1.25), PANEL_BG, GOLD)
    txt(s, "REVM — RISK-ADJUSTED ECONOMIC VALUE METRIC (per shipment)",
        Inches(0.5), Inches(1.3), Inches(8), Inches(0.35), size=10, bold=True, color=GOLD)
    txt(s, "ReVM = contribution_profit  -  risk_penalty  -  time_cost  -  future_impact  -  fx_cost  -  sla_penalty  -  tariff_cost  -  gst_cost",
        Inches(0.5), Inches(1.68), Inches(12.3), Inches(0.68), size=11, color=CYAN)

    # 4 engine blocks
    engines = [
        ("Monte Carlo\nSimulation Engine", [
            "1,000-cycle scenario simulation",
            "Pareto black-swan delays CAPPED at 180 days",
            "(fix: uncapped Pareto produced $500M VaR on $10K cargo)",
            "Outputs: P50/P95 delay, VaR floor, scenario samples",
            "Tier-based SLA penalty rates per contract type",
        ], CYAN, Inches(0.3)),
        ("Cashflow Predictor\nOrchestrator", [
            "Multi-horizon cashflow prediction",
            "Shock detector: flags sudden cashflow anomalies",
            "Root cause engine: traces shock to carrier/route/SKU",
            "Event generator: creates timeline of cash events",
            "Decision support: recommends buffer size adjustments",
        ], GREEN, Inches(3.6)),
        ("Liquidity Score\n& Buffer Engine", [
            "Real-time liquidity scoring (0-100)",
            "Cash buffer requirement calculated per risk tier",
            "Impact engine: models downstream effect of decisions",
            "Concentration engine: flags over-exposure to single carriers",
            "WACC engine: cost of capital updated from FRED API",
        ], PURPLE, Inches(6.9)),
        ("Scenario Simulation\n& Confidence Engine", [
            "What-if scenarios: tariff change, route closure, FX shock",
            "ConfidenceTrustEngine scores model output reliability",
            "Confidence Studio: audit-ready confidence per prediction",
            "FX risk model: USD/INR/EUR/GBP hedging signals",
            "SLA penalty model: tier-based contract breach costs",
        ], AMBER, Inches(10.1)),
    ]

    for title, items, color, bx in engines:
        bw = Inches(3.1)
        card(s, bx, Inches(2.65), bw, Inches(4.55), CARD_BG, color)
        txt(s, title, bx + Inches(0.15), Inches(2.78), bw - Inches(0.3), Inches(0.55),
            size=11, bold=True, color=color)
        for i, item in enumerate(items):
            dot_bullet(s, item, bx + Inches(0.15), Inches(3.42 + i * 0.53),
                       bw - Inches(0.2), dot=color, size=9)

    bottom_bar(s, GOLD)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — RAG / AI COPILOT + AGENTS
# ════════════════════════════════════════════════════════════════════════════
def s08_rag_agents(prs):
    s = blank(prs)
    header_band(s, "RAG AI COPILOT + AUTONOMOUS AGENT SYSTEM",
                "Semantic search over your documents + 5 specialized AI agents dispatched by the orchestrator", PURPLE)

    # Left: RAG
    card(s, Inches(0.3), Inches(1.2), Inches(6.1), Inches(5.9), CARD_BG, PURPLE)
    txt(s, "RAG — Retrieval-Augmented Generation Copilot",
        Inches(0.5), Inches(1.35), Inches(5.7), Inches(0.38), size=13, bold=True, color=PURPLE)

    rag_sections = [
        ("Embedder", "sentence-transformers/all-MiniLM-L6-v2. 384 dimensions, CPU-only (~80ms/batch). "
                     "Graceful fallback to keyword search if library unavailable.", BLUE_LIGHT),
        ("Vector Store", "PostgreSQL + pgvector extension. Embeddings stored alongside document chunks. "
                        "Cosine similarity search over contracts, invoices, financial history.", CYAN),
        ("Ingestion", "PDF, Excel, CSV ingestion via pdfplumber + openpyxl. Chunked, embedded, stored. "
                     "Re-ingestion is idempotent — no duplicate chunks.", GREEN),
        ("Retriever", "Top-K semantic retrieval. Falls back to exact-keyword when embedding model is down. "
                     "Context window managed to fit Claude API limits.", GOLD),
        ("LLM Gateway", "Claude (claude-sonnet) as the reasoning layer. Prompt includes retrieved context + "
                        "user query. Responses grounded in your actual documents.", PURPLE),
    ]
    cy = Inches(1.85)
    for label, desc, color in rag_sections:
        rect(s, Inches(0.5), cy, Inches(0.06), Inches(0.65), color)
        txt(s, label, Inches(0.7), cy, Inches(1.2), Inches(0.28), size=9, bold=True, color=color)
        txt(s, desc, Inches(0.7), cy + Inches(0.3), Inches(5.2), Inches(0.38),
            size=9, color=RGBColor(0xCB, 0xD5, 0xE1))
        cy += Inches(0.85)

    # Right: Agents
    card(s, Inches(6.7), Inches(1.2), Inches(6.33), Inches(5.9), CARD_BG, GOLD)
    txt(s, "5 Autonomous AI Agents",
        Inches(6.9), Inches(1.35), Inches(5.9), Inches(0.38), size=13, bold=True, color=GOLD)

    agents = [
        ("RiskAgent", "Evaluates shipment risk. Pulls XGBoost + GNN scores. "
                     "Triggers reroute recommendation if risk > threshold.", RED_ACC),
        ("FinancialAgent", "Computes ReVM, liquidity score, SLA exposure. "
                          "Flags working capital drain from GST/ITC lag.", GREEN),
        ("RoutingAgent", "Runs GeopoliticalRouteOptimizer. Proposes alternate "
                        "routes when primary is strike-active or HIGH risk.", CYAN),
        ("AnomalyAgent", "Detects statistical anomalies in cashflow, demand, "
                        "and payment patterns using ML drift detection.", PURPLE),
        ("ExecutiveAgent", "Synthesizes all agent outputs into a board-level "
                          "summary. Calls Claude for natural language briefing.", AMBER),
    ]
    ay = Inches(1.9)
    for name, desc, color in agents:
        card(s, Inches(6.9), ay, Inches(5.9), Inches(0.9), PANEL_BG, color)
        txt(s, name, Inches(7.1), ay + Inches(0.08), Inches(2.0), Inches(0.3),
            size=10, bold=True, color=color)
        txt(s, desc, Inches(7.1), ay + Inches(0.42), Inches(5.5), Inches(0.42),
            size=9, color=RGBColor(0xCB, 0xD5, 0xE1))
        ay += Inches(1.05)

    bottom_bar(s, PURPLE)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — ENTERPRISE API + DATA LAYER
# ════════════════════════════════════════════════════════════════════════════
def s09_api(prs):
    s = blank(prs)
    header_band(s, "ENTERPRISE API LAYER + DATA INFRASTRUCTURE",
                "Production FastAPI backend with versioned endpoints, async workers, and multi-DB architecture", BLUE_MID)

    # API endpoints grid
    endpoints = [
        ("/api/v1/predict", "Delay, demand, risk prediction", "POST", PURPLE),
        ("/api/v1/optimize", "Route + inventory optimization", "POST", CYAN),
        ("/api/v1/mapping", "AI field mapping (ERP schema normalize)", "POST", GREEN),
        ("/api/v1/documents", "Document ingestion + RAG queries", "POST", GOLD),
        ("/api/v1/realtime", "WebSocket real-time event stream", "WS", RED_ACC),
        ("/twin/*", "Digital twin simulation endpoints", "GET/POST", BLUE_MID),
        ("/india/*", "GST compliance + refund tracker", "GET/POST", GREEN),
        ("/optimization/*", "MIP optimizer + constraint engine", "POST", AMBER),
        ("/alerts/*", "Alert management + Twilio SMS dispatch", "POST", RED_ACC),
        ("/reports/*", "Financial reports + GSTR generation", "GET", BLUE_LIGHT),
        ("/execution/*", "Dispute engine + reroute executor", "POST", PURPLE),
        ("/confidence/*", "Confidence Studio audit outputs", "GET", CYAN),
    ]

    txt(s, "API ENDPOINTS (15+ routers registered)", Inches(0.35), Inches(1.2),
        Inches(8), Inches(0.35), size=11, bold=True, color=BLUE_LIGHT)
    for i, (ep, desc, method, color) in enumerate(endpoints):
        row, col = i // 3, i % 3
        bx = Inches(0.3 + col * 4.35)
        by = Inches(1.65 + row * 0.7)
        card(s, bx, by, Inches(4.1), Inches(0.58), CARD_BG, color)
        rect(s, bx + Inches(3.35), by + Inches(0.1), Inches(0.62), Inches(0.22), color)
        txt(s, method, bx + Inches(3.35), by + Inches(0.1), Inches(0.62), Inches(0.22),
            size=7, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        txt(s, ep, bx + Inches(0.12), by + Inches(0.06), Inches(3.2), Inches(0.25),
            size=9, bold=True, color=color)
        txt(s, desc, bx + Inches(0.12), by + Inches(0.33), Inches(3.2), Inches(0.22),
            size=8, color=GRAY_MID)

    # Data layer
    card(s, Inches(0.3), Inches(5.55), Inches(12.73), Inches(1.65), PANEL_BG, CYAN)
    txt(s, "DATA & INFRASTRUCTURE LAYER", Inches(0.5), Inches(5.65),
        Inches(6), Inches(0.3), size=10, bold=True, color=CYAN)
    infra = [
        ("PostgreSQL + pgvector", "Primary store + vector embeddings for RAG"),
        ("Redis", "Session cache, FX rate cache, Celery broker"),
        ("Neo4j", "Graph DB for supply chain network topology"),
        ("Celery + Workers", "Async background tasks: model training, report generation"),
        ("SAP / NetSuite Connectors", "ERP write-back + data sync via base connector pattern"),
        ("Docker + Railway", "Containerized deploy, env-var secrets, 1-worker uvicorn"),
    ]
    for i, (label, desc) in enumerate(infra):
        col, row = i % 3, i // 3
        bx = Inches(0.5 + col * 4.2)
        by = Inches(6.05 + row * 0.5)
        txt(s, label + ": ", bx, by, Inches(1.8), Inches(0.38),
            size=9, bold=True, color=GOLD)
        txt(s, desc, bx + Inches(1.85), by, Inches(2.2), Inches(0.38),
            size=9, color=RGBColor(0xCB, 0xD5, 0xE1))

    bottom_bar(s, BLUE_MID)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — CLOSING / AUDIT-READY SUMMARY
# ════════════════════════════════════════════════════════════════════════════
def s10_close(prs):
    s = blank(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    rect(s, 0, 0, Inches(0.12), SLIDE_H, GOLD)
    rect(s, 0, 0, SLIDE_W, Inches(0.08), GOLD)
    rect(s, Inches(7.8), 0, Inches(5.53), SLIDE_H, PANEL_BG)

    txt(s, "WHAT FISCALOGIX GIVES AN AUDITOR",
        Inches(0.35), Inches(0.7), Inches(7.0), Inches(0.55),
        size=22, bold=True, color=WHITE)
    txt(s, "Every decision is traceable. Every cost is quantified. Every risk is explained.",
        Inches(0.35), Inches(1.3), Inches(7.0), Inches(0.38),
        size=12, color=BLUE_LIGHT, italic=True)

    audit_points = [
        (GOLD,    "Full audit trail",        "Every pipeline stage writes to AuditLogger with timestamp, input hash, and output. Immutable."),
        (GREEN,   "SHAP explainability",     "Every risk score comes with top feature importances. No black boxes."),
        (CYAN,    "GST working capital",     "ITC lag and duty drawback quantified per shipment — visible in reports."),
        (PURPLE,  "Monte Carlo VaR",         "1,000-cycle simulation with capped distributions. Credible capital reserve numbers."),
        (AMBER,   "Confidence scores",       "Every ML prediction carries a confidence band. Confidence Studio for drill-down."),
        (RED_ACC, "Decision log",            "What decision was made, by which agent, on which data, at what time — per row."),
    ]
    for i, (color, title, desc) in enumerate(audit_points):
        by = Inches(1.9 + i * 0.82)
        rect(s, Inches(0.35), by + Inches(0.06), Inches(0.32), Inches(0.32), color)
        txt(s, title, Inches(0.8), by, Inches(2.5), Inches(0.38),
            size=11, bold=True, color=color)
        txt(s, desc, Inches(0.8), by + Inches(0.38), Inches(6.7), Inches(0.38),
            size=10, color=RGBColor(0xCB, 0xD5, 0xE1))

    # Right summary panel
    txt(s, "TECH STACK SUMMARY", Inches(8.1), Inches(0.7), Inches(5.0), Inches(0.38),
        size=10, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    stack = [
        ("Language", "Python 3.11"),
        ("API", "FastAPI + Uvicorn"),
        ("ML", "XGBoost, RandomForest, Prophet"),
        ("Deep Learning", "PyTorch GNN"),
        ("AI", "Claude (Anthropic)"),
        ("Embeddings", "all-MiniLM-L6-v2"),
        ("Optimization", "OR-Tools MIP, NetworkX"),
        ("Databases", "PostgreSQL, Redis, Neo4j"),
        ("Queue", "Celery + Redis broker"),
        ("Deploy", "Docker + Railway"),
        ("Compliance", "GST, DGFT, ITC"),
        ("Connectors", "SAP, NetSuite, ERP"),
    ]
    for i, (label, val) in enumerate(stack):
        row, col = i // 2, i % 2
        bx = Inches(8.0 + col * 2.65)
        by = Inches(1.25 + row * 0.52)
        card(s, bx, by, Inches(2.5), Inches(0.42), CARD_BG)
        txt(s, label, bx + Inches(0.1), by + Inches(0.04), Inches(1.15), Inches(0.3),
            size=8, bold=True, color=GRAY_MID)
        txt(s, val, bx + Inches(1.3), by + Inches(0.04), Inches(1.1), Inches(0.3),
            size=8, color=CYAN)

    bottom_bar(s, GOLD)
    rect(s, 0, SLIDE_H - Pt(9), SLIDE_W, Pt(4), RED_ACC)


# ════════════════════════════════════════════════════════════════════════════
# BUILD
# ════════════════════════════════════════════════════════════════════════════
def build():
    prs = new_prs()
    s01_cover(prs)
    s02_why(prs)
    s03_architecture(prs)
    s04_risk(prs)
    s05_gst(prs)
    s06_route(prs)
    s07_financial(prs)
    s08_rag_agents(prs)
    s09_api(prs)
    s10_close(prs)

    out = "fiscalogix_tech_deck.pptx"
    prs.save(out)
    print(f"[OK] Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
